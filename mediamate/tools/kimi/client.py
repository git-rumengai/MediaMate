import asyncio
import os
from typing import Tuple
from mediamate.platforms.parser import XpathParser
from playwright.async_api import Page, BrowserContext, Locator, FrameLocator

from mediamate.config import config
from mediamate.utils.log_manager import log_manager
from mediamate.utils.const import (
    DEFAULT_LOCATION,
    DEFAULT_URL_TIMEOUT,
)
from mediamate.utils.common import (
    get_useragent,
    handle_dialog_accept,
)
from pptx import Presentation
from pptx.util import Cm
from datetime import datetime


logger = log_manager.get_logger(__file__)


class KimiClient:
    kimi_url = 'https://kimi.moonshot.cn/'
    kimi_dir = f'{config.DATA_DIR}/browser/kimi'
    os.makedirs(kimi_dir, exist_ok=True)
    kimi_static = f'{config.PROJECT_DIR}/static'
    kimi_xpath = XpathParser(f'{kimi_static}/elements/kimi/home.yaml')

    async def login(self, playwright, headless: bool = False) -> Tuple[BrowserContext, Page]:
        """ 登录账户 """
        user_data_dir = self.kimi_dir
        # 浏览器启动选项
        launch_options = {
            "user_data_dir": user_data_dir,
            "bypass_csp": False,
            "args": [
                '--disable-blink-features=AutomationControlled',
            ]
        }

        user_agent = get_useragent()
        # 无头模式和有头模式的上下文选项
        context_options_headless = {
            'channel': 'chrome',
            'headless': headless,
            'permissions': ["geolocation"],
            'geolocation': {"longitude": DEFAULT_LOCATION[0], "latitude": DEFAULT_LOCATION[1]},
            'locale': 'zh-CN',
            'timezone_id': "Asia/Shanghai",
            'user_agent': user_agent
        }

        context_options_headed = {
            'channel': 'chrome',
            'headless': headless,
            "viewport": {"width": 1280, "height": 720},
            "slow_mo": 1000,
            'permissions': ["geolocation"],
            'geolocation': {"longitude": DEFAULT_LOCATION[0], "latitude": DEFAULT_LOCATION[1]},
            'locale': 'zh-CN',
            'timezone_id': "Asia/Shanghai",
            'user_agent': user_agent
        }

        if headless:
            logger.info('无头模式启动')
            context = await playwright.chromium.launch_persistent_context(**launch_options, **context_options_headless)
        else:
            logger.info('有头模式启动')
            context = await playwright.chromium.launch_persistent_context(**launch_options, **context_options_headed)

        # 开启跟踪模式
        await context.tracing.start(screenshots=True, snapshots=True)
        page = context.pages[0]
        # 浏览器检测工具
        # await page.goto('https://abrahamjuliot.github.io/creepjs/')
        page.on('dialog', handle_dialog_accept)

        await page.goto(self.kimi_url, timeout=DEFAULT_URL_TIMEOUT)
        await page.wait_for_load_state('load')
        login = await self.get_locator(page, 'ppt login')
        if await login.is_visible():
            logger.info('请登录...')
            await login.click()
            await login.wait_for(timeout=DEFAULT_URL_TIMEOUT * 3, state='hidden')
        return context, page

    async def scroll(self, page: Page):
        """ 直接滚动页面 """
        await page.evaluate("window.scrollBy(0, window.innerHeight)")

    async def wheel(self, page: Page, delta_x: int = 0, delta_y: int = 500):
        """ 模拟滑轮滚动 """
        await page.mouse.wheel(delta_x, delta_y)

    async def ensure_page(self, page) -> Page:
        """ 处理页面加载异常, 只能使用"self.parser.get_xpath" """
        raise NotImplementedError('没有处理可能出现的页面异常')

    async def get_locator(self, page: Page | FrameLocator, parser_path: str) -> Locator:
        """
        以"_"开头是xpath的子节点
        以"_list"结尾是xpath列表
        其他则为一般 xpath 节点
        """
        xpath = self.kimi_xpath.get_xpath(parser_path)
        locator = page.locator(xpath)
        return locator

    async def get_child_locator(self, parent: Locator, child_path: str) -> Locator:
        """ 通过父节点locator定位到子节点 """
        xpath = f'xpath={self.kimi_xpath.get_xpath(child_path)}'
        locator = parent.locator(xpath)
        return locator

    async def get_visible_locator(self, page: Page | FrameLocator, parser_path: str) -> Locator:
        """  """
        locator = await self.get_locator(page, parser_path)
        await locator.wait_for(timeout=DEFAULT_URL_TIMEOUT, state='visible')
        return locator

    async def get_child_visible_locator(self, parent: Locator, child_path: str) -> Locator:
        """ 通过父类locator定位到子类 """
        locator = await self.get_child_locator(parent, child_path)
        await locator.wait_for(timeout=DEFAULT_URL_TIMEOUT, state='visible')
        return locator

    async def get_visible_locators(self, page: Page, parser_path: str) -> Locator:
        """ 传入的是列表的父类, 一定要是唯一的 """
        locator = await self.get_locator(page, parser_path)
        await locator.first.wait_for(timeout=DEFAULT_URL_TIMEOUT, state='visible')
        return locator

    async def get_child_visible_locators(self, parent: Locator, child_path: str) -> Locator:
        """ 通过父类locator定位到子类 """
        locator = await self.get_child_locator(parent, child_path)
        await locator.first.wait_for(timeout=DEFAULT_URL_TIMEOUT, state='visible')
        return locator

    async def wait_first_locator(self, locators: list[Locator]) -> int:
        """ 在给定的 locators 列表中，等待第一个出现的 locator，并返回该 locator 在列表中的索引位置。 """
        tasks = [
            asyncio.create_task(locator.wait_for(state='visible', timeout=DEFAULT_URL_TIMEOUT))
            for locator in locators
        ]
        done, pending = await asyncio.wait(tasks, return_when=asyncio.FIRST_COMPLETED)
        for i, task in enumerate(tasks):
            if task in done and not task.cancelled() and task.exception() is None:
                # 取消所有未完成的任务
                for p in pending:
                    p.cancel()
                return i

        raise TimeoutError('Locator 等待超时')

    def validate_file_path(self, path) -> bool:
        # 检查路径格式
        suffix = path.rsplit('.', 1)
        if len(suffix) > 1 and suffix[1] in ['txt', 'md', 'log', 'doc', 'docx', 'pdf', 'ppt', 'pptx', 'xls', 'xlsx', 'yaml', 'toml']:
            # 检查文件是否存在
            if os.path.exists(path):
                return True
            else:
                raise ValueError(f"文件不存在或非文本文件: '{path}'")
        else:
            return False

    async def download_ppt(self, page: Page, ppt_path: str):
        """ 下载ppt """
        logger.info('正在下载ppt...')
        iframe = page.frame_locator('iframe#aippt-iframe')
        download = await self.get_visible_locator(iframe, 'ppt download')
        await download.click()
        download_button = await self.get_visible_locator(iframe, 'ppt download download_button')
        async with page.expect_download() as download_info:
            await download_button.click()
        download = await download_info.value
        await download.save_as(ppt_path)
        download_finish = await self.get_locator(iframe, 'ppt download download_finish')
        await download_finish.wait_for(timeout=DEFAULT_URL_TIMEOUT * 3, state='visible')
        back_button = await self.get_visible_locator(page, 'ppt back_button')
        await back_button.click()
        logger.info(f'ppt下载完毕: {ppt_path}')

    async def get_ppt(self, page: Page,
                      ppt_path: str,
                      topic: str,
                      scene: str,
                      style: str,
                      color_index: int,
                      card_index: int = 0
                      ) -> Page:
        """  """
        # 1. 进入ppt对话框
        editor = await self.get_visible_locator(page, 'ppt editor')
        await editor.click()
        await page.keyboard.type('@ppt')
        await page.wait_for_timeout(timeout=1)
        scroller_list = await self.get_visible_locators(page, 'ppt scroller_list')
        await scroller_list.nth(0).click()
        await editor.click()
        if self.validate_file_path(topic):
            logger.info('文件转PPT')
            await page.keyboard.type('读取文件内容并生成PPT')
            toolbar = await self.get_visible_locator(page, 'ppt toolbar')
            attach = await self.get_child_visible_locator(toolbar, 'ppt toolbar _attach')
            async with page.expect_file_chooser() as fc_info:
                await attach.click()
            file_chooser = await fc_info.value
            await file_chooser.set_files(topic)
        else:
            logger.info('主题生成 PPT')
            await page.keyboard.type(topic)
        button_click = await self.get_visible_locator(page, 'ppt button_click')
        await button_click.click()

        stop_output = await self.get_visible_locator(page, 'ppt stop_output')
        await stop_output.wait_for(timeout=DEFAULT_URL_TIMEOUT * 5, state='hidden')
        create_ppt = await self.get_locator(page, 'ppt create_ppt')
        edit_ppt = await self.get_locator(page, 'ppt edit_ppt')
        ppt_index = await self.wait_first_locator([create_ppt, edit_ppt])
        if ppt_index == 1:
            logger.info('直接生成ppt')
            await edit_ppt.click()
            await self.download_ppt(page, ppt_path)
        else:
            logger.info('选择模板生成ppt')
            await create_ppt.click()
            # 选择模板
            logger.info(f'选择模板: {scene}, {style}, color: {color_index}, card: {card_index}')
            iframe = page.frame_locator('iframe#aippt-iframe')
            scene_template = iframe.locator(f'//*[text()="{scene}"]')
            style_template = iframe.locator(f'//*[text()="{style}"]')
            colors_list = iframe.locator('//*[text()="主题颜色："]/following-sibling::div/span')
            color_template = colors_list.nth(color_index)
            await scene_template.click()
            await style_template.click()
            await color_template.click()
            cards_list = iframe.locator('//div[@class="template-cover-cards"]/div')
            cards_template = cards_list.nth(card_index)

            await page.wait_for_timeout(timeout=3)
            await cards_template.wait_for(timeout=DEFAULT_URL_TIMEOUT, state='visible')
            await cards_template.click()
            await page.wait_for_timeout(timeout=1)

            new_ppt = await self.get_locator(iframe, 'ppt new_ppt')
            await new_ppt.click()
            finish_text = await self.get_locator(iframe, 'ppt finish_text')
            await finish_text.wait_for(timeout=DEFAULT_URL_TIMEOUT * 5, state='visible')
            to_editor = await self.get_visible_locator(iframe, 'ppt to_editor')
            await to_editor.click()
            page = await self.download_ppt(page, ppt_path)
        return page

    def update_ppt(self, ppt_path: str, username: str, logo_path: str, specified_font: str = 'SimSun') -> str:
        """ 替换ppt中的用户名, 日期和logo """
        presentation = Presentation(ppt_path)
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        horizontal_margin = slide_width / 8
        vertical_margin = slide_height / 6

        current_year = datetime.now().year
        current_date = datetime.now().strftime("%Y.%m")
        for slide_num, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = ''.join([run.text for run in paragraph.runs])
                        if '202X' in paragraph_text and len(paragraph_text.strip()) == 4:
                            paragraph_text = str(current_year)
                        elif '20XX' in paragraph_text and len(paragraph_text.strip()) == 4:
                            paragraph_text = str(current_year)
                        elif '2X' in paragraph_text and len(paragraph_text.strip()) == 2:
                            paragraph_text = str(current_year)[-2:]
                        if slide_num == 0 or slide_num == len(presentation.slides) - 1:
                            if '主讲人' in paragraph_text:
                                paragraph_text = username
                            elif '汇报人' in paragraph_text:
                                paragraph_text = username
                            if "日期" in paragraph_text:
                                paragraph_text = current_date
                            elif '时间' in paragraph_text:
                                paragraph_text = current_date
                            if 'DESIGN' in paragraph_text:
                                slide.shapes._spTree.remove(shape._element)
                        for index, run in enumerate(paragraph.runs):
                            if index == 0:
                                run.text = paragraph_text
                            else:
                                run.text = ''
                            run.font.name = specified_font

                if shape.shape_type == 5 and shape.width / Cm(1) < 5.0 and shape.height / Cm(1) < 0.6:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    if left <= horizontal_margin and top <= vertical_margin:
                        sp = shape._element
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(logo_path, left, top, width, height * 2)
                    elif left >= slide_width - width - horizontal_margin and top <= vertical_margin:
                        sp = shape._element
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(logo_path, left, top, width, height * 2)

        file_name, file_extension = os.path.splitext(os.path.basename(ppt_path))
        update_path = os.path.join(os.path.dirname(ppt_path), f'{file_name}_update{file_extension}')
        presentation.save(update_path)
        return update_path


__all__ = ['KimiClient']
